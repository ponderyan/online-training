#!/usr/bin/env python3
"""
从 PPTX 文件提取全部幻灯片文字，并导出幻灯片内嵌图片（供 OCR 兜底），输出 JSON 到 stdout。
使用方式：
  python3 extract-pptx-text.py <path/to/file.pptx>
输出：
  {"total": N, "slides": [{"index": 1, "text": "..."}, ...], "images": [{"slide": 1, "path": "/tmp/xxx.png"}, ...]}
  - images 仅在幻灯片含 PICTURE 形状时出现（纯图片 PPTX 靠它走 OCR）
  - 图片导出到系统临时目录，调用方 OCR 后负责清理
"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import sys, json, os, tempfile

def extract_pptx_text(pptx_path: str) -> dict:
    prs = Presentation(pptx_path)
    slides = []
    images = []
    tmpdir = None

    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 导出内嵌图片（含 GIF/PNG/JPEG 等），供 OCR 兜底
                try:
                    if tmpdir is None:
                        tmpdir = tempfile.mkdtemp(prefix='pptx-ocr-')
                    ext = shape.image.ext or 'png'
                    img_path = os.path.join(tmpdir, f'slide{i+1}.{ext}')
                    with open(img_path, 'wb') as f:
                        f.write(shape.image.blob)
                    images.append({"slide": i + 1, "path": img_path})
                except Exception:
                    pass
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    t = paragraph.text.strip()
                    if t:
                        texts.append(t)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        t = cell.text.strip()
                        if t:
                            texts.append(t)
            # 也尝试从 组合形状(group shape) 中提取文本
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                try:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            for p in child.text_frame.paragraphs:
                                t = p.text.strip()
                                if t:
                                    texts.append(t)
                except Exception:
                    pass
        slides.append({
            "index": i + 1,
            "text": "\n".join(texts)
        })

    return {"total": len(slides), "slides": slides, "images": images}

if __name__ == '__main__':
    if len(sys.argv) < 2:
        print(json.dumps({"error": "用法: python3 extract-pptx-text.py <pptx_path>"}, ensure_ascii=False))
        sys.exit(1)
    path = sys.argv[1]
    try:
        result = extract_pptx_text(path)
        print(json.dumps(result, ensure_ascii=False))
    except Exception as e:
        print(json.dumps({"error": str(e)}, ensure_ascii=False))
        sys.exit(1)
